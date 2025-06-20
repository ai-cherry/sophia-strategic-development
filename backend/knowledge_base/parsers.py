"""Document Parsers for the Knowledge Base
Handles extracting text content from various document formats like PDF, DOCX, and PPTX.
"""
import logging
from pathlib import Path

import docx
import pandas as pd
import pptx
from pypdf import PdfReader

logger = logging.getLogger(__name__)


class DocumentParser:
    """A class to handle parsing of different document types."""

    @staticmethod
    def parse_pdf(file_path: Path) -> str:
        """Parses a PDF file and extracts its text content.

        Args:
            file_path: The path to the PDF file.

        Returns:
            The extracted text content as a single string.
        """
        logger.info(f"Parsing PDF: {file_path}")
        try:
            reader = PdfReader(file_path)
            text_parts = [
                page.extract_text() for page in reader.pages if page.extract_text()
            ]
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse PDF {file_path}: {e}", exc_info=True)
            return ""

    @staticmethod
    def parse_docx(file_path: Path) -> str:
        """Parses a DOCX file and extracts its text content.

        Args:
            file_path: The path to the DOCX file.

        Returns:
            The extracted text content as a single string.
        """
        logger.info(f"Parsing DOCX: {file_path}")
        try:
            doc = docx.Document(file_path)
            text_parts = [para.text for para in doc.paragraphs if para.text]
            return "\n".join(text_parts)
        except Exception as e:
            logger.error(f"Failed to parse DOCX {file_path}: {e}", exc_info=True)
            return ""

    @staticmethod
    def parse_pptx(file_path: Path) -> str:
        """Parses a PPTX file and extracts its text content from all shapes.

        Args:
            file_path: The path to the PPTX file.

        Returns:
            The extracted text content as a single string.
        """
        logger.info(f"Parsing PPTX: {file_path}")
        try:
            presentation = pptx.Presentation(file_path)
            text_parts = []
            for slide in presentation.slides:
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text.append(shape.text)
                if slide_text:
                    text_parts.append("\n".join(slide_text))
            return "\n\n---\n\n".join(text_parts)  # Separate slides
        except Exception as e:
            logger.error(f"Failed to parse PPTX {file_path}: {e}", exc_info=True)
            return ""

    @staticmethod
    def parse_excel(file_path: Path) -> str:
        """Parses an Excel file (.xlsx, .xls) and converts it to a text-based
        (CSV-like) string for easy embedding.

        Args:
            file_path: The path to the Excel file.

        Returns:
            A string representation of the Excel data.
        """
        logger.info(f"Parsing Excel file: {file_path}")
        try:
            # Read the excel file, assuming the first sheet is the one of interest
            df = pd.read_excel(file_path, engine="openpyxl")
            # Convert the dataframe to a string, which can be chunked and embedded
            # We'll use CSV format for the string representation.
            return df.to_csv(index=False)
        except Exception as e:
            logger.error(f"Failed to parse Excel file {file_path}: {e}", exc_info=True)
            return ""

    @staticmethod
    def parse(file_path: Path) -> str:
        """Automatically detects the file type and parses it.

        Args:
            file_path: The path to the document file.

        Returns:
            The extracted text content.
        """
        suffix = file_path.suffix.lower()
        if suffix == ".pdf":
            return DocumentParser.parse_pdf(file_path)
        elif suffix == ".docx":
            return DocumentParser.parse_docx(file_path)
        elif suffix == ".pptx":
            return DocumentParser.parse_pptx(file_path)
        elif suffix in [".xlsx", ".xls"]:
            return DocumentParser.parse_excel(file_path)
        else:
            logger.warning(f"Unsupported file type: {suffix}. Skipping {file_path}")
            return ""


async def main():
    """A simple main function to test the parsers."""
    # Create dummy files for testing
    from fpdf import FPDF

    test_dir = Path("test_docs")
    test_dir.mkdir(exist_ok=True)

    # Create a dummy PDF
    pdf_path = test_dir / "test.pdf"
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(
        200,
        10,
        txt="This is a test PDF document for Sophia AI's knowledge base.",
        ln=True,
    )
    pdf.output(pdf_path)

    # Create a dummy DOCX
    doc_path = test_dir / "test.docx"
    doc = docx.Document()
    doc.add_paragraph("This is a test DOCX document.")
    doc.add_paragraph("It has multiple paragraphs.")
    doc.save(doc_path)

    # Create a dummy PPTX
    ppt_path = test_dir / "test.pptx"
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = "Test Presentation Slide"
    body_shape = slide.shapes.add_textbox(100, 200, 500, 100)
    body_shape.text_frame.text = "This is the content of the slide."
    prs.save(ppt_path)

    # Test parsers
    pdf_text = DocumentParser.parse(pdf_path)
    print(f"--- Parsed PDF ---\n{pdf_text}\n")

    docx_text = DocumentParser.parse(doc_path)
    print(f"--- Parsed DOCX ---\n{docx_text}\n")

    pptx_text = DocumentParser.parse(ppt_path)
    print(f"--- Parsed PPTX ---\n{pptx_text}\n")

    # Clean up dummy files
    import shutil

    shutil.rmtree(test_dir)


if __name__ == "__main__":
    logging.basicConfig(level=logging.INFO)
    asyncio.run(main())
