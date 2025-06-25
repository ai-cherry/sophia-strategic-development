#!/usr/bin/env python3
"""
Enhanced Ingestion Service for Sophia AI
Handles large files, multiple formats, and provides job-based processing with extensive context windows
"""

import asyncio
import logging
import json
import io
from datetime import datetime, timedelta
from typing import Dict, List, Optional, Any, AsyncIterator
from uuid import uuid4
from pathlib import Path
from enum import Enum
import snowflake.connector
from snowflake.connector import DictCursor
from pydantic import BaseModel
import aiofiles

# Enhanced imports with fallbacks
try:
    import pandas as pd
    HAS_PANDAS = True
except ImportError:
    HAS_PANDAS = False

try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False

try:
    import pptx
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Enhanced Models
class IngestionStatus(str, Enum):
    PENDING = "pending"
    PROCESSING = "processing"
    CHUNKING = "chunking"
    EMBEDDING = "embedding"
    COMPLETED = "completed"
    FAILED = "failed"
    CANCELLED = "cancelled"

class FileType(str, Enum):
    PDF = "application/pdf"
    DOCX = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    DOC = "application/msword"
    TXT = "text/plain"
    CSV = "text/csv"
    JSON = "application/json"
    XLSX = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    XLS = "application/vnd.ms-excel"
    PPTX = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    PPT = "application/vnd.ms-powerpoint"
    MD = "text/markdown"
    RTF = "application/rtf"

class IngestionJob(BaseModel):
    job_id: str
    user_id: str
    filename: str
    file_type: str
    file_size: int
    status: IngestionStatus
    progress: float = 0.0
    chunks_processed: int = 0
    total_chunks: int = 0
    entries_created: List[str] = []
    error_message: Optional[str] = None
    metadata: Dict[str, Any] = {}
    created_at: datetime
    updated_at: datetime
    estimated_completion: Optional[datetime] = None

class DocumentChunk(BaseModel):
    chunk_id: str
    job_id: str
    sequence: int
    content: str
    metadata: Dict[str, Any]
    word_count: int
    character_count: int

class EnhancedIngestionService:
    """Enhanced ingestion service with large file support and contextual processing"""
    
    def __init__(self, snowflake_config: Dict[str, str]):
        self.snowflake_config = snowflake_config
        self.connection = None
        self.active_jobs: Dict[str, IngestionJob] = {}
        
        # Configuration for large file processing
        self.chunk_size = 4000  # Characters per chunk (optimized for context windows)
        self.chunk_overlap = 200  # Overlap between chunks for context preservation
        self.max_file_size = 100 * 1024 * 1024  # 100MB limit
        self.supported_types = {
            FileType.PDF, FileType.DOCX, FileType.DOC, FileType.TXT,
            FileType.CSV, FileType.JSON, FileType.XLSX, FileType.XLS,
            FileType.PPTX, FileType.PPT, FileType.MD, FileType.RTF
        }

    async def connect(self):
        """Connect to Snowflake"""
        try:
            self.connection = snowflake.connector.connect(**self.snowflake_config)
            logger.info("✅ Enhanced Ingestion Service connected to Snowflake")
        except Exception as e:
            logger.error(f"❌ Enhanced Ingestion Service connection failed: {e}")
            raise

    async def disconnect(self):
        """Disconnect from Snowflake"""
        if self.connection:
            self.connection.close()

    async def execute_query(self, query: str, params: Optional[tuple] = None) -> List[Dict[str, Any]]:
        """Execute Snowflake query"""
        try:
            cursor = self.connection.cursor(DictCursor)
            cursor.execute(query, params or ())
            results = cursor.fetchall()
            cursor.close()
            return results
        except Exception as e:
            logger.error(f"Query execution failed: {e}")
            raise

    async def create_ingestion_job(
        self,
        user_id: str,
        filename: str,
        file_content: bytes,
        file_type: str,
        metadata: Optional[Dict[str, Any]] = None
    ) -> IngestionJob:
        """Create a new ingestion job for large file processing"""
        
        job_id = str(uuid4())
        file_size = len(file_content)
        
        # Validate file size
        if file_size > self.max_file_size:
            raise ValueError(f"File size {file_size} exceeds maximum allowed size {self.max_file_size}")
        
        # Validate file type
        if file_type not in [ft.value for ft in self.supported_types]:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        # Estimate processing time based on file size and type
        estimated_completion = self._estimate_completion_time(file_size, file_type)
        
        job = IngestionJob(
            job_id=job_id,
            user_id=user_id,
            filename=filename,
            file_type=file_type,
            file_size=file_size,
            status=IngestionStatus.PENDING,
            metadata=metadata or {},
            created_at=datetime.now(),
            updated_at=datetime.now(),
            estimated_completion=estimated_completion
        )
        
        # Store job in database
        await self._save_job_to_database(job)
        
        # Store in active jobs
        self.active_jobs[job_id] = job
        
        logger.info(f"Created ingestion job {job_id} for file {filename} ({file_size} bytes)")
        return job

    async def process_file_async(self, job_id: str, file_content: bytes) -> None:
        """Process file asynchronously with chunking and context preservation"""
        try:
            job = self.active_jobs.get(job_id)
            if not job:
                logger.error(f"Job {job_id} not found")
                return

            # Update status to processing
            job.status = IngestionStatus.PROCESSING
            job.updated_at = datetime.now()
            await self._update_job_in_database(job)

            # Extract text content based on file type
            text_content = await self._extract_text_content(file_content, job.file_type, job.filename)
            
            # Update status to chunking
            job.status = IngestionStatus.CHUNKING
            await self._update_job_in_database(job)

            # Create intelligent chunks with context preservation
            chunks = await self._create_intelligent_chunks(text_content, job)
            job.total_chunks = len(chunks)
            await self._update_job_in_database(job)

            # Update status to embedding/storage
            job.status = IngestionStatus.EMBEDDING
            await self._update_job_in_database(job)

            # Process chunks and create knowledge entries
            for i, chunk in enumerate(chunks):
                entry_id = await self._process_chunk_to_knowledge_entry(chunk, job)
                job.entries_created.append(entry_id)
                job.chunks_processed = i + 1
                job.progress = (job.chunks_processed / job.total_chunks) * 100
                await self._update_job_in_database(job)
                
                # Small delay to prevent overwhelming the database
                await asyncio.sleep(0.1)

            # Complete the job
            job.status = IngestionStatus.COMPLETED
            job.progress = 100.0
            job.updated_at = datetime.now()
            await self._update_job_in_database(job)

            logger.info(f"Completed ingestion job {job_id}: {len(chunks)} chunks processed")

        except Exception as e:
            logger.error(f"Error processing job {job_id}: {e}")
            if job_id in self.active_jobs:
                job = self.active_jobs[job_id]
                job.status = IngestionStatus.FAILED
                job.error_message = str(e)
                job.updated_at = datetime.now()
                await self._update_job_in_database(job)

    async def _extract_text_content(self, file_content: bytes, file_type: str, filename: str) -> str:
        """Extract text content from various file types with enhanced support"""
        try:
            if file_type == FileType.TXT.value or file_type == FileType.MD.value:
                return file_content.decode('utf-8', errors='ignore')
            
            elif file_type == FileType.CSV.value and HAS_PANDAS:
                try:
                    df = pd.read_csv(io.BytesIO(file_content))
                    text = f"CSV Data from {filename}:\n\n"
                    text += f"Columns: {', '.join(df.columns)}\n"
                    text += f"Total Records: {len(df)}\n\n"
                    
                    # Include all data for context, not just sample
                    for index, row in df.iterrows():
                        row_text = " | ".join([f"{col}: {str(val)}" for col, val in row.items()])
                        text += f"Row {index + 1}: {row_text}\n"
                    
                    return text
                except Exception:
                    return file_content.decode('utf-8', errors='ignore')
            
            elif file_type == FileType.JSON.value:
                data = json.loads(file_content.decode('utf-8'))
                return f"JSON Data from {filename}:\n\n{json.dumps(data, indent=2)}"
            
            elif file_type == FileType.PDF.value and HAS_PDF:
                try:
                    pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_content))
                    text = f"PDF Document: {filename}\n\n"
                    for page_num, page in enumerate(pdf_reader.pages):
                        page_text = page.extract_text()
                        if page_text.strip():
                            text += f"Page {page_num + 1}:\n{page_text}\n\n"
                    return text
                except Exception as e:
                    logger.warning(f"PDF extraction failed: {e}")
                    return f"PDF Document: {filename}\n(Text extraction failed)"
            
            elif file_type == FileType.DOCX.value and HAS_DOCX:
                try:
                    doc = DocxDocument(io.BytesIO(file_content))
                    text = f"Word Document: {filename}\n\n"
                    for paragraph in doc.paragraphs:
                        if paragraph.text.strip():
                            text += paragraph.text + "\n"
                    return text
                except Exception as e:
                    logger.warning(f"DOCX extraction failed: {e}")
                    return f"Word Document: {filename}\n(Text extraction failed)"
            
            elif file_type in [FileType.XLSX.value, FileType.XLS.value] and HAS_EXCEL:
                try:
                    # Handle Excel files
                    if file_type == FileType.XLSX.value:
                        workbook = openpyxl.load_workbook(io.BytesIO(file_content))
                    else:
                        # For XLS, try with pandas
                        df = pd.read_excel(io.BytesIO(file_content))
                        text = f"Excel Data from {filename}:\n\n"
                        text += f"Columns: {', '.join(df.columns)}\n"
                        text += f"Total Records: {len(df)}\n\n"
                        text += df.to_string(index=False)
                        return text
                    
                    # Process XLSX with openpyxl
                    text = f"Excel Workbook: {filename}\n\n"
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text += f"Sheet: {sheet_name}\n"
                        for row in sheet.iter_rows(values_only=True):
                            if any(cell is not None for cell in row):
                                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                                text += f"{row_text}\n"
                        text += "\n"
                    return text
                except Exception as e:
                    logger.warning(f"Excel extraction failed: {e}")
                    return f"Excel Document: {filename}\n(Text extraction failed)"
            
            elif file_type in [FileType.PPTX.value, FileType.PPT.value] and HAS_PPTX:
                try:
                    presentation = pptx.Presentation(io.BytesIO(file_content))
                    text = f"PowerPoint Presentation: {filename}\n\n"
                    for slide_num, slide in enumerate(presentation.slides):
                        text += f"Slide {slide_num + 1}:\n"
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                text += f"{shape.text}\n"
                        text += "\n"
                    return text
                except Exception as e:
                    logger.warning(f"PowerPoint extraction failed: {e}")
                    return f"PowerPoint Presentation: {filename}\n(Text extraction failed)"
            
            else:
                # Fallback to raw text extraction
                return file_content.decode('utf-8', errors='ignore')
                
        except Exception as e:
            logger.error(f"Text extraction failed for {filename}: {e}")
            return f"Content from {filename} (extraction failed: {str(e)})"

    async def _create_intelligent_chunks(self, text_content: str, job: IngestionJob) -> List[DocumentChunk]:
        """Create intelligent chunks with context preservation and large context windows"""
        chunks = []
        
        # Split text into sentences for better chunking
        sentences = self._split_into_sentences(text_content)
        
        current_chunk = ""
        current_metadata = {
            "filename": job.filename,
            "file_type": job.file_type,
            "job_id": job.job_id,
            "user_id": job.user_id
        }
        
        sequence = 0
        
        for sentence in sentences:
            # Check if adding this sentence would exceed chunk size
            if len(current_chunk) + len(sentence) > self.chunk_size and current_chunk:
                # Create chunk with overlap from previous chunk for context
                chunk = DocumentChunk(
                    chunk_id=str(uuid4()),
                    job_id=job.job_id,
                    sequence=sequence,
                    content=current_chunk.strip(),
                    metadata=current_metadata.copy(),
                    word_count=len(current_chunk.split()),
                    character_count=len(current_chunk)
                )
                chunks.append(chunk)
                
                # Start new chunk with overlap for context preservation
                overlap_start = max(0, len(current_chunk) - self.chunk_overlap)
                current_chunk = current_chunk[overlap_start:] + " " + sentence
                sequence += 1
            else:
                current_chunk += " " + sentence if current_chunk else sentence
        
        # Add final chunk if there's remaining content
        if current_chunk.strip():
            chunk = DocumentChunk(
                chunk_id=str(uuid4()),
                job_id=job.job_id,
                sequence=sequence,
                content=current_chunk.strip(),
                metadata=current_metadata.copy(),
                word_count=len(current_chunk.split()),
                character_count=len(current_chunk)
            )
            chunks.append(chunk)
        
        logger.info(f"Created {len(chunks)} intelligent chunks for job {job.job_id}")
        return chunks

    def _split_into_sentences(self, text: str) -> List[str]:
        """Split text into sentences for better chunking"""
        import re
        
        # Simple sentence splitting - can be enhanced with NLP libraries
        sentences = re.split(r'(?<=[.!?])\s+', text)
        return [s.strip() for s in sentences if s.strip()]

    async def _process_chunk_to_knowledge_entry(self, chunk: DocumentChunk, job: IngestionJob) -> str:
        """Process chunk into knowledge entry with enhanced context"""
        entry_id = str(uuid4())
        
        # Create enhanced title with context
        title = f"{job.filename} - Part {chunk.sequence + 1}"
        if chunk.sequence == 0:
            title = f"{job.filename} - Introduction"
        
        # Enhanced content with metadata
        content = f"Source: {job.filename}\n"
        content += f"File Type: {job.file_type}\n"
        content += f"Part: {chunk.sequence + 1}\n"
        content += f"Word Count: {chunk.word_count}\n\n"
        content += chunk.content
        
        # Auto-categorize based on filename and content
        category_id = self._auto_categorize_content(job.filename, chunk.content)
        
        # Ensure category exists
        await self._ensure_category_exists(category_id)
        
        # Enhanced metadata
        metadata = {
            **chunk.metadata,
            "chunk_id": chunk.chunk_id,
            "sequence": chunk.sequence,
            "word_count": chunk.word_count,
            "character_count": chunk.character_count,
            "processing_timestamp": datetime.now().isoformat(),
            "large_file_processing": True
        }
        
        # Insert knowledge entry
        insert_query = """
        INSERT INTO KNOWLEDGE_BASE_ENTRIES 
        (ENTRY_ID, TITLE, CONTENT, CATEGORY_ID, STATUS, METADATA, CREATED_AT, UPDATED_AT)
        VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
        """
        
        now = datetime.now()
        await self.execute_query(insert_query, (
            entry_id, title, content, category_id, "published",
            json.dumps(metadata), now, now
        ))
        
        return entry_id

    def _auto_categorize_content(self, filename: str, content: str) -> str:
        """Auto-categorize content based on filename and content analysis"""
        filename_lower = filename.lower()
        content_lower = content.lower()
        
        # Customer-related keywords
        if any(keyword in filename_lower or keyword in content_lower 
               for keyword in ['customer', 'client', 'contact', 'company', 'tenant']):
            return "customers"
        
        # Product-related keywords
        elif any(keyword in filename_lower or keyword in content_lower 
                for keyword in ['product', 'service', 'feature', 'solution', 'offering']):
            return "products"
        
        # Employee-related keywords
        elif any(keyword in filename_lower or keyword in content_lower 
                for keyword in ['employee', 'staff', 'team', 'personnel', 'hr']):
            return "employees"
        
        # Financial keywords
        elif any(keyword in filename_lower or keyword in content_lower 
                for keyword in ['financial', 'revenue', 'budget', 'accounting']):
            return "financial"
        
        # Process keywords
        elif any(keyword in filename_lower or keyword in content_lower 
                for keyword in ['process', 'procedure', 'workflow', 'manual']):
            return "processes"
        
        else:
            return "general"

    async def _ensure_category_exists(self, category_id: str):
        """Ensure category exists in database"""
        check_query = "SELECT CATEGORY_ID FROM KNOWLEDGE_CATEGORIES WHERE CATEGORY_ID = %s"
        existing = await self.execute_query(check_query, (category_id,))
        
        if not existing:
            category_names = {
                "customers": "Customer Information",
                "products": "Products & Services", 
                "employees": "Employee Directory",
                "financial": "Financial Information",
                "processes": "Business Processes",
                "general": "General Knowledge"
            }
            
            category_name = category_names.get(category_id, category_id.title())
            
            create_query = """
            INSERT INTO KNOWLEDGE_CATEGORIES (CATEGORY_ID, CATEGORY_NAME, DESCRIPTION, CREATED_AT)
            VALUES (%s, %s, %s, %s)
            """
            
            await self.execute_query(create_query, (
                category_id, category_name, f"Auto-created category", datetime.now()
            ))

    def _estimate_completion_time(self, file_size: int, file_type: str) -> datetime:
        """Estimate completion time based on file size and type"""
        # Base processing time (seconds per MB)
        base_time_per_mb = {
            FileType.TXT.value: 2,
            FileType.CSV.value: 5,
            FileType.JSON.value: 3,
            FileType.PDF.value: 10,
            FileType.DOCX.value: 8,
            FileType.XLSX.value: 15,
            FileType.PPTX.value: 12
        }
        
        time_per_mb = base_time_per_mb.get(file_type, 5)
        file_size_mb = file_size / (1024 * 1024)
        estimated_seconds = file_size_mb * time_per_mb
        
        return datetime.now() + timedelta(seconds=estimated_seconds)

    async def _save_job_to_database(self, job: IngestionJob):
        """Save ingestion job to database"""
        query = """
        INSERT INTO INGESTION_JOBS 
        (JOB_ID, USER_ID, FILENAME, FILE_TYPE, FILE_SIZE, STATUS, PROGRESS, 
         CHUNKS_PROCESSED, TOTAL_CHUNKS, METADATA, CREATED_AT, UPDATED_AT, ESTIMATED_COMPLETION)
        VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s, %s)
        """
        
        # Ensure the table exists
        await self._ensure_ingestion_tables_exist()
        
        await self.execute_query(query, (
            job.job_id, job.user_id, job.filename, job.file_type, job.file_size,
            job.status.value, job.progress, job.chunks_processed, job.total_chunks,
            json.dumps(job.metadata), job.created_at, job.updated_at, job.estimated_completion
        ))

    async def _update_job_in_database(self, job: IngestionJob):
        """Update ingestion job in database"""
        query = """
        UPDATE INGESTION_JOBS SET 
        STATUS = %s, PROGRESS = %s, CHUNKS_PROCESSED = %s, TOTAL_CHUNKS = %s,
        UPDATED_AT = %s, ERROR_MESSAGE = %s
        WHERE JOB_ID = %s
        """
        
        await self.execute_query(query, (
            job.status.value, job.progress, job.chunks_processed, job.total_chunks,
            job.updated_at, job.error_message, job.job_id
        ))

    async def _ensure_ingestion_tables_exist(self):
        """Ensure ingestion tracking tables exist"""
        create_jobs_table = """
        CREATE TABLE IF NOT EXISTS INGESTION_JOBS (
            JOB_ID VARCHAR(255) PRIMARY KEY,
            USER_ID VARCHAR(255) NOT NULL,
            FILENAME VARCHAR(500) NOT NULL,
            FILE_TYPE VARCHAR(100) NOT NULL,
            FILE_SIZE INTEGER NOT NULL,
            STATUS VARCHAR(50) NOT NULL,
            PROGRESS FLOAT DEFAULT 0.0,
            CHUNKS_PROCESSED INTEGER DEFAULT 0,
            TOTAL_CHUNKS INTEGER DEFAULT 0,
            ERROR_MESSAGE TEXT,
            METADATA VARIANT,
            CREATED_AT TIMESTAMP_NTZ NOT NULL,
            UPDATED_AT TIMESTAMP_NTZ NOT NULL,
            ESTIMATED_COMPLETION TIMESTAMP_NTZ
        )
        """
        
        await self.execute_query(create_jobs_table)

    async def get_job_status(self, job_id: str) -> Optional[IngestionJob]:
        """Get ingestion job status"""
        if job_id in self.active_jobs:
            return self.active_jobs[job_id]
        
        # Query from database
        query = "SELECT * FROM INGESTION_JOBS WHERE JOB_ID = %s"
        results = await self.execute_query(query, (job_id,))
        
        if results:
            row = results[0]
            return IngestionJob(
                job_id=row['JOB_ID'],
                user_id=row['USER_ID'],
                filename=row['FILENAME'],
                file_type=row['FILE_TYPE'],
                file_size=row['FILE_SIZE'],
                status=IngestionStatus(row['STATUS']),
                progress=row['PROGRESS'],
                chunks_processed=row['CHUNKS_PROCESSED'],
                total_chunks=row['TOTAL_CHUNKS'],
                error_message=row['ERROR_MESSAGE'],
                metadata=json.loads(row['METADATA']) if row['METADATA'] else {},
                created_at=row['CREATED_AT'],
                updated_at=row['UPDATED_AT'],
                estimated_completion=row['ESTIMATED_COMPLETION']
            )
        
        return None

    async def get_recent_jobs(self, user_id: str, limit: int = 20) -> List[IngestionJob]:
        """Get recent ingestion jobs for user"""
        query = """
        SELECT * FROM INGESTION_JOBS 
        WHERE USER_ID = %s 
        ORDER BY CREATED_AT DESC 
        LIMIT %s
        """
        
        results = await self.execute_query(query, (user_id, limit))
        
        jobs = []
        for row in results:
            job = IngestionJob(
                job_id=row['JOB_ID'],
                user_id=row['USER_ID'],
                filename=row['FILENAME'],
                file_type=row['FILE_TYPE'],
                file_size=row['FILE_SIZE'],
                status=IngestionStatus(row['STATUS']),
                progress=row['PROGRESS'],
                chunks_processed=row['CHUNKS_PROCESSED'],
                total_chunks=row['TOTAL_CHUNKS'],
                error_message=row['ERROR_MESSAGE'],
                metadata=json.loads(row['METADATA']) if row['METADATA'] else {},
                created_at=row['CREATED_AT'],
                updated_at=row['UPDATED_AT'],
                estimated_completion=row['ESTIMATED_COMPLETION']
            )
            jobs.append(job)
        
        return jobs

    async def cancel_job(self, job_id: str) -> bool:
        """Cancel an ingestion job"""
        if job_id in self.active_jobs:
            job = self.active_jobs[job_id]
            job.status = IngestionStatus.CANCELLED
            job.updated_at = datetime.now()
            await self._update_job_in_database(job)
            del self.active_jobs[job_id]
            return True
        return False 